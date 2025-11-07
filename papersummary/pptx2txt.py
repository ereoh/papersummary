"""
Convert ppts text to txt
"""

from papersummary.base import BaseTextExtractor

from pptx import Presentation


class PPTX2TextExtractor(BaseTextExtractor):

    def __init__(
        self,
        default_prompt: str = "Write a clear, concise, objective summary for the following document:",
    ):
        super().__init__(supported_extensions=[".pptx"], default_prompt=default_prompt)

    def _extract_text(self, file: str) -> str:
        # Extract text from the .pptx file

        presentation = Presentation(file)

        text = []
        for i, slide in enumerate(presentation.slides):
            text.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text_frame = shape.text_frame

                for paragraph in text_frame.paragraphs:
                    # Append the text, stripping whitespace
                    if paragraph.text:
                        text.append(paragraph.text.strip())

        return "\n".join(text)
