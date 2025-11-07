"""
Convert ppts text to txt
"""

from papersummary.base import BaseTextExtractor
from papersummary.utils import DEFAULT_PROMPT

from pptx import Presentation


class PPTX2TextExtractor(BaseTextExtractor):
    """Microsoft .pptx text extractor 
    """
    def __init__(
        self,
        default_prompt: str = DEFAULT_PROMPT,
    ):
        """Initializes .pptx text extractor

        Args:
            default_prompt (str, optional): Summary prompt to use. Defaults to DEFAULT_PROMPT.
        """ 
        super().__init__(supported_extensions=[".pptx"], default_prompt=default_prompt)

    def _extract_text(self, file: str) -> str:
        """Extracts text from .pptx file

        Args:
            file (str): Path to .pptx file

        Returns:
            str: Extracted text.
        """

        presentation = Presentation(file)

        text = []
        for i, slide in enumerate(presentation.slides):
            text.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text_frame = shape.text_frame

                for paragraph in text_frame.paragraphs:
                    # Append the text, stripping whitespace
                    if paragraph.text:
                        text.append(paragraph.text.strip())

        return "\n".join(text)
